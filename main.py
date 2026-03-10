"""Telegram bot powered by Google Gemini AI.

Features:
- General conversation via Gemini
- /presentation <topic> – generate and send a PPTX presentation on any topic
"""

import io
import json
import logging
import os
import re
import textwrap

import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from telegram import Update
from telegram.ext import (
    Application,
    CommandHandler,
    ContextTypes,
    MessageHandler,
    filters,
)

load_dotenv()

logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO,
)
logger = logging.getLogger(__name__)

TELEGRAM_TOKEN = os.environ["TELEGRAM_TOKEN"]
GEMINI_API_KEY = os.environ["GEMINI_API_KEY"]

genai.configure(api_key=GEMINI_API_KEY)
_model = genai.GenerativeModel("gemini-1.5-flash")

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x16, 0x21, 0x3E)
HIGHLIGHT = RGBColor(0x0F, 0x3D, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)
ORANGE = RGBColor(0xFF, 0x8F, 0x00)

TELEGRAM_MSG_LIMIT = 4096

# RJ-45 T568B wire colours (name, RGB)
RJ45_COLORS = [
    ("Бело-оранжевый / White-Orange", RGBColor(0xFF, 0xCC, 0x99)),
    ("Оранжевый / Orange", RGBColor(0xFF, 0x80, 0x00)),
    ("Бело-зелёный / White-Green", RGBColor(0xCC, 0xFF, 0xCC)),
    ("Синий / Blue", RGBColor(0x00, 0x66, 0xFF)),
    ("Бело-синий / White-Blue", RGBColor(0x99, 0xBB, 0xFF)),
    ("Зелёный / Green", RGBColor(0x00, 0xAA, 0x44)),
    ("Бело-коричневый / White-Brown", RGBColor(0xDD, 0xBB, 0x99)),
    ("Коричневый / Brown", RGBColor(0x88, 0x44, 0x00)),
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, color: RGBColor) -> None:
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = WHITE,
    align=PP_ALIGN.LEFT,
) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK_BG)

    # Decorative bar at top
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(10), Inches(0.12),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BLUE
    bar.line.fill.background()

    # Bottom bar
    bar2 = slide.shapes.add_shape(
        1, Inches(0), Inches(7.38), Inches(10), Inches(0.12)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = LIGHT_BLUE
    bar2.line.fill.background()

    _add_textbox(
        slide, title,
        left=0.5, top=2.5, width=9, height=1.5,
        font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, subtitle,
        left=0.5, top=4.2, width=9, height=1,
        font_size=20, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER,
    )


def _add_content_slide(
    prs: Presentation,
    heading: str,
    bullets: list[str],
    slide_num: int,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_color = DARK_BG if slide_num % 2 == 0 else ACCENT
    _set_slide_bg(slide, bg_color)

    # Heading bar
    hdr = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = HIGHLIGHT
    hdr.line.fill.background()

    _add_textbox(
        slide, heading,
        left=0.3, top=0.1, width=9.4, height=0.9,
        font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
    )

    # Bullets
    y = 1.3
    for bullet in bullets:
        _add_textbox(
            slide, f"▸  {bullet}",
            left=0.5, top=y, width=9, height=0.55,
            font_size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        )
        y += 0.6


def _add_rj45_slide(prs: Presentation) -> None:
    """Dedicated slide showing T568B RJ-45 crimping colour order."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)

    hdr = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = HIGHLIGHT
    hdr.line.fill.background()

    _add_textbox(
        slide, "Схема обжима RJ-45  (стандарт T568B)",
        left=0.3, top=0.1, width=9.4, height=0.9,
        font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
    )

    # Draw coloured wire boxes
    box_w, box_h = Inches(0.9), Inches(0.6)
    start_x = Inches(0.35)
    y_top = Inches(1.3)

    for i, (name, rgb) in enumerate(RJ45_COLORS):
        x = start_x + i * Inches(1.1)

        # Pin number
        pin_tb = slide.shapes.add_textbox(x, y_top, box_w, Inches(0.3))
        p = pin_tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = LIGHT_BLUE

        # Coloured rectangle
        rect = slide.shapes.add_shape(1, x, y_top + Inches(0.35), box_w, box_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = rgb
        rect.line.color.rgb = WHITE

        # Wire label (rotated text via textbox below)
        label_tb = slide.shapes.add_textbox(
            x - Inches(0.1), y_top + Inches(1.05), box_w + Inches(0.2), Inches(0.9)
        )
        label_tb.text_frame.word_wrap = True
        p2 = label_tb.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = name
        run2.font.size = Pt(9)
        run2.font.color.rgb = WHITE

    # T568A note
    _add_textbox(
        slide,
        "T568A: пары 2 и 3 меняются местами  "
        "(1-Бело-зелёный, 2-Зелёный, 3-Бело-оранжевый, 6-Оранжевый)",
        left=0.35, top=6.1, width=9.3, height=0.7,
        font_size=12, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.LEFT,
    )

    _add_textbox(
        slide,
        "Важно: оба конца кабеля должны быть обжаты по одному стандарту.",
        left=0.35, top=6.8, width=9.3, height=0.5,
        font_size=12, bold=False, color=ORANGE, align=PP_ALIGN.LEFT,
    )


def _add_closing_slide(prs: Presentation, topic: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)

    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BLUE
    bar.line.fill.background()

    _add_textbox(
        slide, "Спасибо за внимание!",
        left=0.5, top=2.8, width=9, height=1.0,
        font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, topic,
        left=0.5, top=4.0, width=9, height=0.7,
        font_size=20, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Gemini helpers
# ---------------------------------------------------------------------------

def _ask_gemini(prompt: str) -> str:
    response = _model.generate_content(prompt)
    return response.text


def _generate_presentation_outline(topic: str) -> list[dict]:
    """Ask Gemini to produce a JSON outline for the presentation."""
    prompt = textwrap.dedent(f"""
        Ты помощник, создающий профессиональные презентации.
        Тема: «{topic}»
        Создай структуру презентации: ровно 7 содержательных слайдов.
        Для каждого слайда дай заголовок и 4-5 кратких тезисов (не длиннее 12 слов каждый).
        Ответь ТОЛЬКО корректным JSON-массивом без каких-либо пояснений.
        Формат:
        [
          {{"heading": "...", "bullets": ["...", "...", "...", "...", "..."]}},
          ...
        ]
    """)
    raw = _ask_gemini(prompt).strip()
    # Strip markdown code fences if present
    raw = re.sub(r"^```[a-z]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw)
    return json.loads(raw)


def _build_pptx(topic: str, outline: list[dict]) -> bytes:
    """Build the PPTX file and return it as bytes."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, topic, "Информационные технологии и сети")

    for i, slide_data in enumerate(outline, start=1):
        _add_content_slide(prs, slide_data["heading"], slide_data["bullets"], i)

    # RJ-45 crimping slide
    _add_rj45_slide(prs)

    _add_closing_slide(prs, topic)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ---------------------------------------------------------------------------
# Telegram handlers
# ---------------------------------------------------------------------------

async def start_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    await update.message.reply_text(
        "👋 Привет! Я бот на базе Google Gemini.\n\n"
        "Команды:\n"
        "  /presentation <тема> — создать презентацию (PPTX)\n"
        "  /help — справка\n\n"
        "Или просто напишите мне что-нибудь, и я отвечу."
    )


async def help_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    await update.message.reply_text(
        "ℹ️ Справка:\n\n"
        "  /presentation <тема>\n"
        "    Генерирует и присылает PPTX-презентацию по заданной теме.\n"
        "    Пример: /presentation среда передачи данных\n\n"
        "  Свободный чат:\n"
        "    Напишите любой вопрос — Gemini ответит."
    )


async def presentation_handler(
    update: Update, context: ContextTypes.DEFAULT_TYPE
) -> None:
    topic = " ".join(context.args).strip() if context.args else ""
    if not topic:
        await update.message.reply_text(
            "⚠️ Укажите тему. Пример:\n/presentation среда передачи данных"
        )
        return

    msg = await update.message.reply_text(
        f"⏳ Генерирую презентацию на тему «{topic}»…"
    )
    try:
        outline = _generate_presentation_outline(topic)
        pptx_bytes = _build_pptx(topic, outline)
        safe_name = re.sub(r"[^\w\-]", "_", topic)[:50]
        filename = f"{safe_name}.pptx"
        await update.message.reply_document(
            document=io.BytesIO(pptx_bytes),
            filename=filename,
            caption=f"📊 Презентация: «{topic}»",
        )
        await msg.delete()
    except (json.JSONDecodeError, ValueError, OSError) as exc:
        logger.exception("Error generating presentation")
        await msg.edit_text(f"❌ Не удалось создать презентацию: {exc}")


async def chat_handler(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_text = update.message.text or ""
    if not user_text:
        return
    try:
        reply = _ask_gemini(user_text)
        for chunk in [reply[i: i + TELEGRAM_MSG_LIMIT] for i in range(0, len(reply), TELEGRAM_MSG_LIMIT)]:
            await update.message.reply_text(chunk)
    except (ValueError, OSError) as exc:
        logger.exception("Gemini error")
        await update.message.reply_text(f"❌ Ошибка Gemini: {exc}")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main() -> None:
    app = Application.builder().token(TELEGRAM_TOKEN).build()
    app.add_handler(CommandHandler("start", start_handler))
    app.add_handler(CommandHandler("help", help_handler))
    app.add_handler(CommandHandler("presentation", presentation_handler))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, chat_handler))
    logger.info("Bot started")
    app.run_polling()


if __name__ == "__main__":
    main()
